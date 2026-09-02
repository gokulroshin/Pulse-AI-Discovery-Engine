"""Generate Discovery Engine Workflow Word Document and 2-Slide PowerPoint."""

import os
import sys

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu as PEmu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════════════
# STYLING HELPERS
# ═══════════════════════════════════════════════════════════════════

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_CELL = RGBColor(0x1E, 0x29, 0x3B)


def set_cell_bg(cell, color_hex):
    """Set background color of a table cell."""
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def add_styled_table(doc, headers, rows, col_widths=None):
    """Add a professionally styled table."""
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"

    # Header row
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run(h)
        run.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_cell_bg(cell, "1E293B")

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.rows[r_idx + 1].cells[c_idx]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.size = Pt(8.5)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            bg = "F1F5F9" if r_idx % 2 == 0 else "FFFFFF"
            set_cell_bg(cell, bg)

    if col_widths:
        for i, w in enumerate(col_widths):
            for row in table.rows:
                row.cells[i].width = Inches(w)

    return table


def add_heading_with_color(doc, text, level, color=None):
    """Add a heading with optional color."""
    h = doc.add_heading(text, level=level)
    if color:
        for run in h.runs:
            run.font.color.rgb = color
    return h


# ═══════════════════════════════════════════════════════════════════
# WORD DOCUMENT
# ═══════════════════════════════════════════════════════════════════

def generate_word_doc(output_path):
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10)
    style.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    # ── TITLE PAGE ──
    for _ in range(5):
        doc.add_paragraph("")

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("PULSE AI DISCOVERY ENGINE")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0x63, 0x66, 0xF1)
    r.font.name = "Calibri"

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = subtitle.add_run("Detailed System Workflow & Architecture")
    r2.font.size = Pt(16)
    r2.font.color.rgb = LIGHT_GRAY
    r2.font.name = "Calibri"

    line = doc.add_paragraph()
    line.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = line.add_run("━" * 50)
    r3.font.color.rgb = RGBColor(0x63, 0x66, 0xF1)
    r3.font.size = Pt(10)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r4 = meta.add_run("Myntra Graduation Project  •  August 2026\nMulti-Channel Consumer Intelligence Platform")
    r4.font.size = Pt(11)
    r4.font.color.rgb = LIGHT_GRAY
    r4.font.name = "Calibri"

    doc.add_page_break()

    # ── TABLE OF CONTENTS ──
    add_heading_with_color(doc, "Table of Contents", 1, RGBColor(0x63, 0x66, 0xF1))
    toc_items = [
        "1. Executive Overview",
        "2. System Architecture (4-Layer Pipeline)",
        "3. Layer 1 — Multi-Channel Ingestion",
        "4. Layer 2 — LLM Signal Extraction",
        "5. Layer 3 — Semantic Aggregation & Scoring",
        "6. Layer 4 — Intelligence Surface (Dashboard & AI Q&A)",
        "7. Data Model & Database Schema",
        "8. API Endpoints Reference",
        "9. Technology Stack",
        "10. Current Corpus Statistics",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item, style="List Number")
        p.paragraph_format.space_after = Pt(2)

    doc.add_page_break()

    # ── 1. EXECUTIVE OVERVIEW ──
    add_heading_with_color(doc, "1. Executive Overview", 1, RGBColor(0x63, 0x66, 0xF1))

    doc.add_paragraph(
        "The Pulse AI Discovery Engine is a 4-layer qualitative research pipeline that transforms "
        "unstructured consumer feedback from multiple independent channels (Reddit, Google Play Store, "
        "Apple App Store, YouTube, and manual uploads) into ranked, evidence-grounded business opportunity "
        "areas for fashion e-commerce."
    )
    doc.add_paragraph(
        "The core business problem it addresses is Myntra's 30-day wishlist-to-purchase conversion gap — "
        "understanding why shoppers save fashion items to their wishlist but never buy them. "
        "The engine processes 1,938 raw documents into 1,554 verified qualitative signals, clusters them "
        "into 8 unique opportunity areas, and ranks them using a 5-dimension composite scoring formula."
    )

    add_heading_with_color(doc, "Key Capabilities", 2, RGBColor(0x8B, 0x5C, 0xF6))
    capabilities = [
        ("Multi-Channel Scraping", "Automated ingestion from 5 independent consumer feedback sources with SHA-256 deduplication."),
        ("LLM-Powered Extraction", "Google Gemini Flash extracts structured behavioral signals with a verbatim quote hallucination guardrail."),
        ("Semantic Clustering", "TF-IDF + SVD embeddings fed into agglomerative clustering with silhouette-optimized k selection."),
        ("5-Dimension Scoring", "Composite ranking using Frequency, Triangulation, Conversion Relevance, Segment Breadth, and Actionability."),
        ("AI-Powered Q&A", "RAG-based question answering with topic-matched corroborating customer evidence from across channels."),
        ("Interactive Dashboard", "Next.js 14 dark-mode dashboard with opportunity drill-down, evidence explorer, and segment analytics."),
    ]
    for title_text, desc in capabilities:
        p = doc.add_paragraph()
        r = p.add_run(f"• {title_text}: ")
        r.bold = True
        r.font.size = Pt(10)
        p.add_run(desc).font.size = Pt(10)

    doc.add_page_break()

    # ── 2. SYSTEM ARCHITECTURE ──
    add_heading_with_color(doc, "2. System Architecture (4-Layer Pipeline)", 1, RGBColor(0x63, 0x66, 0xF1))

    doc.add_paragraph(
        "The engine operates as a sequential 4-layer pipeline. Each layer processes the output of the "
        "previous layer and adds a higher level of intelligence:"
    )

    add_styled_table(doc,
        ["Layer", "Name", "Input", "Output", "Key Technology"],
        [
            ["Layer 1", "Multi-Channel Ingestion", "Raw reviews, posts, comments from 5 platforms", "1,938 normalized, deduplicated, enriched documents in SQLite", "Python scrapers, SHA-256 hashing, regex metadata enrichment"],
            ["Layer 2", "LLM Signal Extraction", "1,938 raw documents", "1,554 verified behavioral extractions (reason + quote + signal type)", "Google Gemini Flash, structured JSON, verbatim quote guardrail"],
            ["Layer 3", "Aggregation & Scoring", "1,554 discrete extractions", "8 unique ranked opportunity areas with composite scores", "TF-IDF/SVD, Agglomerative Clustering, 5-dimension scoring"],
            ["Layer 4", "Intelligence Surface", "Ranked opportunities + corpus", "Interactive dashboard + AI Q&A with evidence", "Next.js 14, FastAPI, Gemini RAG synthesis"],
        ],
        col_widths=[0.7, 1.3, 1.6, 1.8, 1.8],
    )

    doc.add_page_break()

    # ── 3. LAYER 1 ──
    add_heading_with_color(doc, "3. Layer 1 — Multi-Channel Ingestion", 1, RGBColor(0x63, 0x66, 0xF1))

    doc.add_paragraph(
        "The ingestion pipeline scrapes, normalizes, deduplicates, and enriches raw consumer feedback "
        "from 5 independent data sources into a unified corpus."
    )

    add_heading_with_color(doc, "Source Channels", 2, RGBColor(0x06, 0xB6, 0xD4))
    add_styled_table(doc,
        ["Channel", "Scraper Module", "Data Type", "Engagement Metric"],
        [
            ["Google Play Store", "PlayStoreScraper", "App reviews & ratings", "Helpful votes"],
            ["Apple App Store", "AppStoreScraper", "App reviews & ratings", "Helpful votes"],
            ["Reddit", "RedditScraper", "Posts & comments from fashion subreddits", "Upvotes"],
            ["YouTube", "YouTubeScraper", "Video comments on try-on/haul videos", "Likes"],
            ["Manual Upload", "ManualUploadHandler", "CSV/JSON files from research teams", "Custom score"],
        ],
        col_widths=[1.3, 1.5, 2.0, 1.2],
    )

    add_heading_with_color(doc, "Processing Steps", 2, RGBColor(0x06, 0xB6, 0xD4))
    steps = [
        ("Step 1 — Scraping", "Each scraper module fetches up to 200 documents per source. Documents are returned as RawScrapedDocument objects with content text, platform, URL, and engagement score."),
        ("Step 2 — Text Normalization", "The TextNormalizer cleans raw text by unescaping HTML entities, stripping HTML tags, removing non-printable characters, collapsing whitespace, and truncating to 4,000 characters. Documents shorter than 15 characters or with fewer than 10 alphanumeric characters are discarded."),
        ("Step 3 — SHA-256 Deduplication", "A deterministic SHA-256 hash is computed on the canonicalized (lowercased, whitespace-collapsed) text. Duplicates are detected both within the current batch and against all existing documents in the database."),
        ("Step 4 — Metadata Enrichment", "The MetadataEnricher classifies each document along 3 heuristic dimensions using keyword pattern matching: Product Category (ethnic_wear, western, footwear, accessories, general), Gender Context (women, men, unisex, unknown), and Brand Tier (premium, mid, value, unknown)."),
        ("Step 5 — Persistence", "All new, unique, enriched documents are bulk-inserted into the raw_documents SQLite table with a foreign key to their PipelineRun record for full traceability."),
    ]
    for title_text, desc in steps:
        p = doc.add_paragraph()
        r = p.add_run(f"{title_text}: ")
        r.bold = True
        r.font.size = Pt(10)
        p.add_run(desc).font.size = Pt(10)
        p.paragraph_format.space_after = Pt(6)

    doc.add_page_break()

    # ── 4. LAYER 2 ──
    add_heading_with_color(doc, "4. Layer 2 — LLM Signal Extraction", 1, RGBColor(0x63, 0x66, 0xF1))

    doc.add_paragraph(
        "This layer transforms raw, unstructured consumer text into discrete, typed behavioral signals — "
        "each with a synthesized reason statement, a validated verbatim quote, a confidence level, and a "
        "signal type classification."
    )

    add_heading_with_color(doc, "Extraction Schema", 2, RGBColor(0x8B, 0x5C, 0xF6))
    add_styled_table(doc,
        ["Field", "Type", "Description"],
        [
            ["reason_text", "string", "1-sentence synthesis of the specific friction, behavior, or motivation"],
            ["verbatim_quote", "string", "Exact substring from source text (validated against original)"],
            ["confidence", "high / medium / low", "Explicitly stated → high, clearly implied → medium, inferred → low"],
            ["signal_type", "friction / motivation / behavior / uncertainty / comparison / external_validation", "Behavioral signal classification"],
            ["preliminary_cluster_hint", "string (optional)", "2-4 word topic tag for clustering assistance (e.g., sizing_doubt)"],
        ],
        col_widths=[1.5, 1.8, 3.5],
    )

    add_heading_with_color(doc, "Processing Steps", 2, RGBColor(0x8B, 0x5C, 0xF6))
    extraction_steps = [
        ("Checkpoint-Resilient Fetching", "Only documents not yet extracted are queried, enabling safe restarts if the pipeline is interrupted mid-batch."),
        ("Batch Chunking", "Documents are split into batches of 20 for efficient Gemini API calls."),
        ("Gemini Flash API Call", "Each batch is sent to Google Gemini Flash with a system instruction prompt, few-shot examples, and structured JSON output schema. Temperature is set to 0.1 for deterministic extraction."),
        ("Token Bucket Rate Limiting", "A rate limiter (60 requests/minute, burst limit of 5) prevents API 429 errors."),
        ("Verbatim Quote Hallucination Guardrail", "CRITICAL: Every verbatim_quote returned by the LLM is checked as an exact substring match against the original source document text. Both strings are canonicalized (lowercased, whitespace-collapsed). Quotes that fail this check are flagged as hallucinated and silently discarded. This ensures 100% of stored evidence is traceable to real customer words."),
        ("Database Persistence", "Valid extractions are bulk-inserted with foreign keys to their source document and pipeline run."),
    ]
    for title_text, desc in extraction_steps:
        p = doc.add_paragraph()
        r = p.add_run(f"• {title_text}: ")
        r.bold = True
        r.font.size = Pt(10)
        p.add_run(desc).font.size = Pt(10)
        p.paragraph_format.space_after = Pt(6)

    doc.add_page_break()

    # ── 5. LAYER 3 ──
    add_heading_with_color(doc, "5. Layer 3 — Semantic Aggregation & Scoring", 1, RGBColor(0x63, 0x66, 0xF1))

    doc.add_paragraph(
        "This layer groups the 1,554 discrete extractions into coherent thematic opportunity areas, "
        "assigns human-readable labels, and computes a multi-dimensional composite business relevance "
        "score for each."
    )

    add_heading_with_color(doc, "4-Step Aggregation Coordinator", 2, RGBColor(0x10, 0xB9, 0x81))

    agg_steps = [
        ("Step 1/4 — Generate Embeddings", "Each extraction's reason_text is vectorized using a TF-IDF + Truncated SVD pipeline (scikit-learn). This produces a 100-dimensional dense unit-normalized embedding for each extraction. An in-memory cache prevents redundant computation."),
        ("Step 2/4 — Agglomerative Clustering", "The embeddings are clustered using hierarchical agglomerative clustering with cosine distance and average linkage. The optimal number of clusters (k) is automatically determined by searching k=8 to k=15 and selecting the k with the highest silhouette score. The current corpus yields k=15 raw clusters with a silhouette score of 0.342."),
        ("Step 3/4 — Taxonomy Label Synthesis & Consolidation", "Each cluster's top-5 exemplar reasons and quotes are matched against 12 predefined domain keyword dictionaries to assign a human-readable label and description. Clusters that receive the same label are then CONSOLIDATED: their extraction IDs are merged and representative quotes are deduplicated. This reduces 15 raw clusters down to 8 unique opportunity areas."),
        ("Step 4/4 — Composite Opportunity Scoring", "Each unique taxonomy node is scored across 5 business-relevance dimensions and ranked by a weighted composite score."),
    ]
    for title_text, desc in agg_steps:
        p = doc.add_paragraph()
        r = p.add_run(f"{title_text}: ")
        r.bold = True
        r.font.size = Pt(10)
        p.add_run(desc).font.size = Pt(10)
        p.paragraph_format.space_after = Pt(8)

    add_heading_with_color(doc, "Composite Scoring Formula", 2, RGBColor(0x10, 0xB9, 0x81))
    doc.add_paragraph(
        "Composite Score = 0.25 × Frequency + 0.25 × Triangulation + 0.25 × Conversion Relevance "
        "+ 0.15 × Segment Breadth + 0.10 × Actionability"
    ).runs[0].bold = True

    add_styled_table(doc,
        ["Dimension", "Weight", "Calculation", "What It Measures"],
        [
            ["Frequency", "25%", "extraction_count / max_extractions", "How prevalent is this pain point in the corpus?"],
            ["Triangulation", "25%", "distinct_platforms / 4 (capped at 1.0)", "Is this confirmed across independent sources?"],
            ["Conversion Relevance", "25%", "Domain-expert heuristic matrix", "How directly does this block wishlist-to-purchase?"],
            ["Segment Breadth", "15%", "Normalized Shannon entropy", "Does this affect a niche or the broad customer base?"],
            ["Actionability", "10%", "Non-monetary feasibility heuristic", "Can it be fixed without discounting?"],
        ],
        col_widths=[1.2, 0.6, 1.8, 2.6],
    )

    add_heading_with_color(doc, "Ranked Opportunity Areas (Current Output)", 2, RGBColor(0x10, 0xB9, 0x81))
    add_styled_table(doc,
        ["Rank", "Opportunity Area", "Score", "Extractions", "Confidence"],
        [
            ["#1", "Styling & Outfit Context Deficit", "0.87", "642", "HIGH"],
            ["#2", "Post-Order & Return Policy Friction", "0.64", "370", "HIGH"],
            ["#3", "Fit & Sizing Confidence Gap", "0.62", "198", "HIGH"],
            ["#4", "Wishlist Decision Deferral & Intent Latency", "0.51", "21", "HIGH"],
            ["#5", "Post-Confirmation Inventory & Order Cancellations", "0.50", "77", "HIGH"],
            ["#6", "Review Authenticity & Trust Deficit", "0.45", "141", "HIGH"],
            ["#7", "Bookmarking vs. High-Intent Ambiguity", "0.39", "76", "HIGH"],
            ["#8", "Fulfillment & Delivery Tracking Friction", "0.36", "29", "HIGH"],
        ],
        col_widths=[0.5, 2.8, 0.6, 0.9, 0.8],
    )

    doc.add_page_break()

    # ── 6. LAYER 4 ──
    add_heading_with_color(doc, "6. Layer 4 — Intelligence Surface", 1, RGBColor(0x63, 0x66, 0xF1))

    add_heading_with_color(doc, "6a. Interactive Dashboard (Next.js 14)", 2, RGBColor(0xF5, 0x9E, 0x0B))
    dashboard_features = [
        ("Corpus Stats Panel", "Displays total documents (1,938), verified extractions (1,554), opportunity areas (8), and triangulation channels (4)."),
        ("Ranked Opportunity Table", "Sortable by composite score, triangulation, frequency. Each row shows the opportunity label, description, score bar, and source channel badges."),
        ("Opportunity Detail View", "Drill-down view with paginated verbatim evidence explorer, per-platform filtering, and representative customer quotes."),
        ("Segment Analysis", "Category, gender, and brand tier breakdowns with proportional bar charts."),
    ]
    for title_text, desc in dashboard_features:
        p = doc.add_paragraph()
        r = p.add_run(f"• {title_text}: ")
        r.bold = True
        p.add_run(desc)

    add_heading_with_color(doc, "6b. AI Insight Search (RAG Q&A Engine)", 2, RGBColor(0xF5, 0x9E, 0x0B))
    doc.add_paragraph(
        "The AI Insight Search allows product managers to ask natural language questions about consumer "
        "behavior and receive evidence-grounded answers with corroborating customer quotes from across channels."
    )

    qa_steps = [
        ("Question Intent Detection", "The engine first checks if the question matches one of 3 pre-curated Core Knowledge Templates (why_wishlist, purchase_prevention, uncertainties_remaining). For custom questions, it extracts content keywords and maps them to domain synonyms."),
        ("Targeted Evidence Retrieval", "The retrieve_supporting_evidence() function scores candidate quotes by keyword relevance (+4 per keyword match in quote, +2 in reason), penalizes off-topic quotes (-15 for unrelated app bugs), and enforces multi-channel diversity (one quote per platform first)."),
        ("Gemini RAG Synthesis", "For dynamic questions, evidence and opportunity themes are formatted into a context prompt and sent to Gemini Flash for structured JSON synthesis in simple, everyday language."),
        ("Offline Fallback", "If Gemini is unavailable, the engine constructs a contextual response from the top-scoring evidence quotes and opportunity descriptions."),
    ]
    for title_text, desc in qa_steps:
        p = doc.add_paragraph()
        r = p.add_run(f"• {title_text}: ")
        r.bold = True
        p.add_run(desc)
        p.paragraph_format.space_after = Pt(6)

    doc.add_page_break()

    # ── 7. DATA MODEL ──
    add_heading_with_color(doc, "7. Data Model & Database Schema", 1, RGBColor(0x63, 0x66, 0xF1))

    add_styled_table(doc,
        ["Table", "Purpose", "Key Fields", "Relationships"],
        [
            ["raw_documents", "Stores normalized, enriched consumer feedback", "doc_id (PK), source_platform, content_text, content_hash (unique), engagement_score, inferred_category, inferred_gender_context, inferred_brand_tier", "→ extractions (1:many), → pipeline_runs"],
            ["extractions", "Discrete behavioral signals extracted by LLM", "extraction_id (PK), doc_id (FK), reason_text, verbatim_quote, confidence, signal_type, taxonomy_node_id (FK)", "→ raw_documents (many:1), → taxonomy_nodes (many:1)"],
            ["taxonomy_nodes", "Thematic opportunity area clusters", "node_id (PK), label, description, extraction_count, representative_quotes (JSON), parent_node_id (FK, self-ref)", "→ extractions (1:many), → opportunity_scores (1:many)"],
            ["opportunity_scores", "Multi-dimensional business scores per opportunity", "score_id (PK), taxonomy_node_id (FK), frequency_score, triangulation_score, conversion_relevance_score, segment_breadth_score, actionability_score, composite_score, rank", "→ taxonomy_nodes (many:1)"],
            ["pipeline_runs", "Execution tracking for ingestion, extraction, scoring", "run_id (PK), stage, status, config (JSON), stats (JSON), created_at, completed_at", "→ raw_documents, extractions, opportunity_scores"],
        ],
        col_widths=[1.1, 1.3, 2.5, 1.8],
    )

    doc.add_page_break()

    # ── 8. API ENDPOINTS ──
    add_heading_with_color(doc, "8. API Endpoints Reference", 1, RGBColor(0x63, 0x66, 0xF1))

    add_styled_table(doc,
        ["Method", "Endpoint", "Purpose"],
        [
            ["GET", "/health", "Health check and server status"],
            ["GET", "/api/v1/corpus/stats", "Corpus statistics (doc count, extraction count, platform breakdown)"],
            ["GET", "/api/v1/opportunities", "Ranked opportunity areas with composite scores"],
            ["GET", "/api/v1/opportunities/{id}", "Single opportunity detail with full evidence"],
            ["GET", "/api/v1/opportunities/{id}/evidence", "Paginated verbatim evidence for an opportunity"],
            ["POST", "/api/v1/insights/ask", "AI Q&A with RAG synthesis and corroborating evidence"],
            ["GET", "/api/v1/segments/{dim}/breakdown", "Segment breakdown (category, gender, brand_tier)"],
            ["POST", "/api/v1/pipeline/run", "Trigger ingestion + extraction pipeline"],
            ["GET", "/api/v1/taxonomy", "Taxonomy node hierarchy"],
        ],
        col_widths=[0.7, 2.5, 3.0],
    )

    # ── 9. TECH STACK ──
    add_heading_with_color(doc, "9. Technology Stack", 1, RGBColor(0x63, 0x66, 0xF1))

    add_styled_table(doc,
        ["Component", "Technology", "Purpose"],
        [
            ["Backend API", "FastAPI (Python 3.10+)", "RESTful API server with async support"],
            ["LLM Engine", "Google Gemini Flash 3.5", "Structured extraction + RAG synthesis"],
            ["Database", "SQLite (SQLAlchemy ORM)", "Persistent storage for corpus, extractions, taxonomy, scores"],
            ["Embeddings", "TF-IDF + Truncated SVD (scikit-learn)", "Dense vector representations for semantic clustering"],
            ["Clustering", "Agglomerative Clustering (scikit-learn)", "Hierarchical semantic grouping with silhouette optimization"],
            ["Frontend", "Next.js 14 (React + TypeScript)", "Interactive dark-mode dashboard with real-time API calls"],
            ["Styling", "Custom CSS (glassmorphism theme)", "Premium visual design system"],
        ],
        col_widths=[1.2, 2.3, 3.0],
    )

    # ── 10. STATS ──
    add_heading_with_color(doc, "10. Current Corpus Statistics", 1, RGBColor(0x63, 0x66, 0xF1))

    add_styled_table(doc,
        ["Metric", "Value"],
        [
            ["Raw Documents Ingested", "1,938"],
            ["Verified Qualitative Extractions", "1,554"],
            ["Unique Opportunity Areas", "8"],
            ["Triangulation Channels", "4 (Reddit, Play Store, App Store, YouTube)"],
            ["Top Opportunity", "Styling & Outfit Context Deficit (Score: 0.87)"],
            ["Aggregation Pipeline Duration", "~1.3 seconds"],
            ["Hallucination Guardrail Rejection Rate", "Tracked per pipeline run"],
        ],
        col_widths=[2.5, 4.0],
    )

    doc.save(output_path)
    print(f"Word document saved: {output_path}")


# ═══════════════════════════════════════════════════════════════════
# POWERPOINT (2 SLIDES)
# ═══════════════════════════════════════════════════════════════════

def generate_pptx(output_path):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # ── SLIDE 1: Architecture Overview ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    bg = slide1.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRGBColor(0x0F, 0x17, 0x2A)

    # Title
    title_box = slide1.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12), PInches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "PULSE AI DISCOVERY ENGINE — System Architecture"
    r.font.size = PPt(32)
    r.font.bold = True
    r.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Calibri"

    # Subtitle
    sub_box = slide1.shapes.add_textbox(PInches(0.5), PInches(1.0), PInches(12), PInches(0.5))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Multi-Channel Consumer Intelligence Platform  •  Myntra 30-Day Wishlist-to-Purchase Conversion"
    r2.font.size = PPt(14)
    r2.font.color.rgb = PRGBColor(0x94, 0xA3, 0xB8)
    r2.font.name = "Calibri"

    # 4 Layer Boxes
    layers = [
        {
            "title": "LAYER 1",
            "name": "Multi-Channel Ingestion",
            "desc": "• 5 Sources: Play Store, App Store,\n  Reddit, YouTube, Manual Upload\n• Text normalization & SHA-256 dedup\n• Metadata enrichment (category,\n  gender, brand tier)\n• 1,938 documents ingested",
            "color": PRGBColor(0x06, 0xB6, 0xD4),
            "bg": PRGBColor(0x0E, 0x2A, 0x3A),
            "x": 0.4,
        },
        {
            "title": "LAYER 2",
            "name": "LLM Signal Extraction",
            "desc": "• Google Gemini Flash structured\n  JSON extraction\n• Verbatim quote hallucination\n  guardrail (substring validation)\n• 6 signal types: friction, motivation,\n  behavior, uncertainty, comparison\n• 1,554 verified extractions",
            "color": PRGBColor(0x8B, 0x5C, 0xF6),
            "bg": PRGBColor(0x1E, 0x15, 0x3A),
            "x": 3.5,
        },
        {
            "title": "LAYER 3",
            "name": "Aggregation & Scoring",
            "desc": "• TF-IDF + SVD embeddings\n• Agglomerative clustering (k=15)\n• 12-theme taxonomy labeling\n• Label consolidation → 8 unique areas\n• 5-dimension composite scoring:\n  Freq + Tri + Conv + Seg + Act\n• Ranked #1 through #8",
            "color": PRGBColor(0x10, 0xB9, 0x81),
            "bg": PRGBColor(0x0A, 0x2A, 0x1E),
            "x": 6.6,
        },
        {
            "title": "LAYER 4",
            "name": "Intelligence Surface",
            "desc": "• Next.js 14 interactive dashboard\n• Ranked opportunity table\n• Evidence drill-down explorer\n• AI Insight Search (RAG Q&A)\n• Topic-matched corroborating\n  customer evidence\n• Segment analytics",
            "color": PRGBColor(0xF5, 0x9E, 0x0B),
            "bg": PRGBColor(0x2A, 0x20, 0x0A),
            "x": 9.7,
        },
    ]

    for layer in layers:
        x = PInches(layer["x"])
        y = PInches(1.7)
        w = PInches(2.9)
        h = PInches(4.5)

        # Card background
        shape = slide1.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = layer["bg"]
        shape.line.color.rgb = layer["color"]
        shape.line.width = PEmu(18000)
        shape.shadow.inherit = False

        # Layer number
        num_box = slide1.shapes.add_textbox(x + PInches(0.15), y + PInches(0.15), PInches(2.6), PInches(0.35))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = layer["title"]
        nr.font.size = PPt(11)
        nr.font.bold = True
        nr.font.color.rgb = layer["color"]
        nr.font.name = "Calibri"

        # Layer name
        name_box = slide1.shapes.add_textbox(x + PInches(0.15), y + PInches(0.45), PInches(2.6), PInches(0.55))
        nmtf = name_box.text_frame
        nmp = nmtf.paragraphs[0]
        nmr = nmp.add_run()
        nmr.text = layer["name"]
        nmr.font.size = PPt(16)
        nmr.font.bold = True
        nmr.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
        nmr.font.name = "Calibri"

        # Description
        desc_box = slide1.shapes.add_textbox(x + PInches(0.15), y + PInches(1.1), PInches(2.6), PInches(3.2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = layer["desc"]
        dr.font.size = PPt(11)
        dr.font.color.rgb = PRGBColor(0xCB, 0xD5, 0xE1)
        dr.font.name = "Calibri"

    # Arrows between layers
    for i in range(3):
        x_start = PInches(layers[i]["x"] + 3.0)
        y_arrow = PInches(3.8)
        arrow_box = slide1.shapes.add_textbox(x_start, y_arrow, PInches(0.5), PInches(0.5))
        atf = arrow_box.text_frame
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = "→"
        ar.font.size = PPt(28)
        ar.font.color.rgb = PRGBColor(0x63, 0x66, 0xF1)
        ar.font.bold = True

    # Bottom stats bar
    stats_items = [
        ("1,938", "Documents"),
        ("1,554", "Extractions"),
        ("8", "Opportunity Areas"),
        ("4", "Channels"),
        ("0.87", "Top Score"),
    ]
    for i, (val, label) in enumerate(stats_items):
        sx = PInches(0.5 + i * 2.5)
        sy = PInches(6.5)

        stat_box = slide1.shapes.add_textbox(sx, sy, PInches(2.2), PInches(0.9))
        stf = stat_box.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = val
        sr.font.size = PPt(22)
        sr.font.bold = True
        sr.font.color.rgb = PRGBColor(0x63, 0x66, 0xF1)
        sr.font.name = "Calibri"

        sp2 = stf.add_paragraph()
        sp2.alignment = PP_ALIGN.CENTER
        sr2 = sp2.add_run()
        sr2.text = label
        sr2.font.size = PPt(10)
        sr2.font.color.rgb = PRGBColor(0x94, 0xA3, 0xB8)
        sr2.font.name = "Calibri"

    # ── SLIDE 2: Ranked Opportunities & Scoring ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg2 = slide2.background
    fill2 = bg2.fill
    fill2.solid()
    fill2.fore_color.rgb = PRGBColor(0x0F, 0x17, 0x2A)

    # Title
    t2_box = slide2.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12), PInches(0.7))
    t2f = t2_box.text_frame
    t2p = t2f.paragraphs[0]
    t2r = t2p.add_run()
    t2r.text = "Ranked Business Opportunity Areas & Composite Scoring"
    t2r.font.size = PPt(30)
    t2r.font.bold = True
    t2r.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
    t2r.font.name = "Calibri"

    # Scoring formula box
    formula_box = slide2.shapes.add_textbox(PInches(0.5), PInches(1.1), PInches(12.5), PInches(0.6))
    fmtf = formula_box.text_frame
    fmp = fmtf.paragraphs[0]
    fmp.alignment = PP_ALIGN.LEFT
    fmr = fmp.add_run()
    fmr.text = "Composite Score  =  25% Frequency  +  25% Triangulation  +  25% Conversion Relevance  +  15% Segment Breadth  +  10% Actionability"
    fmr.font.size = PPt(13)
    fmr.font.bold = True
    fmr.font.color.rgb = PRGBColor(0x10, 0xB9, 0x81)
    fmr.font.name = "Calibri"

    # Opportunity rankings
    opps = [
        ("#1", "Styling & Outfit Context Deficit", "0.87", "642", "Users struggle to visualize outfits with existing wardrobe", PRGBColor(0x10, 0xB9, 0x81)),
        ("#2", "Post-Order & Return Policy Friction", "0.64", "370", "Fear of difficult returns, delayed refunds, denied exchanges", PRGBColor(0x06, 0xB6, 0xD4)),
        ("#3", "Fit & Sizing Confidence Gap", "0.62", "198", "Size charts don't match real body; inconsistent across brands", PRGBColor(0x8B, 0x5C, 0xF6)),
        ("#4", "Wishlist Decision Deferral & Intent Latency", "0.51", "21", "Saving items as preliminary step; waiting for sales/opinions", PRGBColor(0xF5, 0x9E, 0x0B)),
        ("#5", "Post-Confirmation Inventory & Order Cancellations", "0.50", "77", "Confirmed orders cancelled due to stock sync failures", PRGBColor(0xF5, 0x9E, 0x0B)),
        ("#6", "Review Authenticity & Trust Deficit", "0.45", "141", "Distrust in generic 5-star reviews without real customer photos", PRGBColor(0xEF, 0x44, 0x44)),
        ("#7", "Bookmarking vs. High-Intent Ambiguity", "0.39", "76", "Wishlist used as moodboard/inspiration, not purchase funnel", PRGBColor(0xEF, 0x44, 0x44)),
        ("#8", "Fulfillment & Delivery Tracking Friction", "0.36", "29", "Delayed shipments, fake delivery attempts, no tracking updates", PRGBColor(0xEF, 0x44, 0x44)),
    ]

    start_y = 1.85
    row_h = 0.62

    for i, (rank, name, score, count, desc, color) in enumerate(opps):
        y = PInches(start_y + i * row_h)

        # Rank badge
        rank_box = slide2.shapes.add_textbox(PInches(0.5), y, PInches(0.5), PInches(0.5))
        rtf = rank_box.text_frame
        rtf.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = rtf.paragraphs[0].add_run()
        rr.text = rank
        rr.font.size = PPt(14)
        rr.font.bold = True
        rr.font.color.rgb = color
        rr.font.name = "Calibri"

        # Name
        name_box = slide2.shapes.add_textbox(PInches(1.1), y, PInches(3.5), PInches(0.5))
        ntf = name_box.text_frame
        nr = ntf.paragraphs[0].add_run()
        nr.text = name
        nr.font.size = PPt(13)
        nr.font.bold = True
        nr.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
        nr.font.name = "Calibri"

        # Score
        score_box = slide2.shapes.add_textbox(PInches(4.7), y, PInches(0.7), PInches(0.5))
        stf = score_box.text_frame
        stf.paragraphs[0].alignment = PP_ALIGN.CENTER
        sr = stf.paragraphs[0].add_run()
        sr.text = score
        sr.font.size = PPt(14)
        sr.font.bold = True
        sr.font.color.rgb = color
        sr.font.name = "Calibri"

        # Count
        cnt_box = slide2.shapes.add_textbox(PInches(5.5), y, PInches(0.9), PInches(0.5))
        ctf = cnt_box.text_frame
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cr = ctf.paragraphs[0].add_run()
        cr.text = f"{count} quotes"
        cr.font.size = PPt(10)
        cr.font.color.rgb = PRGBColor(0x94, 0xA3, 0xB8)
        cr.font.name = "Calibri"

        # Score bar
        bar_width = float(score) / 1.0 * 2.8
        bar = slide2.shapes.add_shape(1, PInches(6.6), y + PEmu(120000), PInches(bar_width), PEmu(140000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Description
        desc_box = slide2.shapes.add_textbox(PInches(9.6), y, PInches(3.5), PInches(0.5))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dr = dtf.paragraphs[0].add_run()
        dr.text = desc
        dr.font.size = PPt(9.5)
        dr.font.color.rgb = PRGBColor(0x94, 0xA3, 0xB8)
        dr.font.name = "Calibri"

    # Column headers
    headers = [
        (0.5, "RANK"), (1.1, "OPPORTUNITY AREA"), (4.7, "SCORE"),
        (5.5, "EVIDENCE"), (6.6, "COMPOSITE BAR"), (9.6, "DESCRIPTION"),
    ]
    for hx, htxt in headers:
        hbox = slide2.shapes.add_textbox(PInches(hx), PInches(1.55), PInches(1.5), PInches(0.3))
        htf = hbox.text_frame
        hr = htf.paragraphs[0].add_run()
        hr.text = htxt
        hr.font.size = PPt(8)
        hr.font.bold = True
        hr.font.color.rgb = PRGBColor(0x63, 0x66, 0xF1)
        hr.font.name = "Calibri"

    # Bottom: Scoring dimension legend
    dims = [
        ("Frequency (25%)", "Prevalence across corpus"),
        ("Triangulation (25%)", "Cross-platform confirmation"),
        ("Conversion Rel. (25%)", "Impact on wishlist→purchase"),
        ("Segment Breadth (15%)", "Affects broad or niche users"),
        ("Actionability (10%)", "Non-monetary fix feasibility"),
    ]
    for i, (dim_name, dim_desc) in enumerate(dims):
        dx = PInches(0.4 + i * 2.5)
        dy = PInches(6.6)
        dbox = slide2.shapes.add_textbox(dx, dy, PInches(2.3), PInches(0.8))
        dtf = dbox.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dr1 = dp.add_run()
        dr1.text = dim_name
        dr1.font.size = PPt(10)
        dr1.font.bold = True
        dr1.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
        dr1.font.name = "Calibri"
        dp2 = dtf.add_paragraph()
        dr2 = dp2.add_run()
        dr2.text = dim_desc
        dr2.font.size = PPt(8)
        dr2.font.color.rgb = PRGBColor(0x94, 0xA3, 0xB8)
        dr2.font.name = "Calibri"

    prs.save(output_path)
    print(f"PowerPoint saved: {output_path}")


# ═══════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    docx_path = os.path.join(base_dir, "Pulse_AI_Discovery_Engine_Workflow.docx")
    pptx_path = os.path.join(base_dir, "Pulse_AI_Discovery_Engine_Overview.pptx")

    generate_word_doc(docx_path)
    generate_pptx(pptx_path)

    print(f"\nFiles generated:")
    print(f"  Word:  {docx_path}")
    print(f"  PPTX:  {pptx_path}")
